from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


# Widescreen 16:9 — mọi export (kể cả có template) dùng cùng hệ tọa độ
REF_W = 13.333
REF_H = 7.5
REF_MARGIN = 0.55


@dataclass
class SlideCanvas:
    w: float = REF_W
    h: float = REF_H
    margin: float = REF_MARGIN
    template_mode: bool = False

    @property
    def content_w(self) -> float:
        return self.w - 2 * self.margin


# Active canvas for current export (set in build_pptx)
CV = SlideCanvas()

THEMES = {
    "professional": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title": RGBColor(0x0F, 0x17, 0x2A),
        "body": RGBColor(0x33, 0x41, 0x55),
        "muted": RGBColor(0x64, 0x74, 0x8B),
        "accent": RGBColor(0x4F, 0x46, 0xE5),
        "accent_soft": RGBColor(0xEE, 0xF2, 0xFF),
        "card": RGBColor(0xF8, 0xFA, 0xFC),
        "card_border": RGBColor(0xE2, 0xE8, 0xF0),
        "on_accent": RGBColor(0xFF, 0xFF, 0xFF),
    },
    "dark": {
        "bg": RGBColor(0x0F, 0x17, 0x2A),
        "title": RGBColor(0xF8, 0xFA, 0xFC),
        "body": RGBColor(0xCB, 0xD5, 0xE1),
        "muted": RGBColor(0x94, 0xA3, 0xB8),
        "accent": RGBColor(0x38, 0xBD, 0xF8),
        "accent_soft": RGBColor(0x1E, 0x29, 0x3B),
        "card": RGBColor(0x1E, 0x29, 0x3B),
        "card_border": RGBColor(0x33, 0x41, 0x55),
        "on_accent": RGBColor(0x0F, 0x17, 0x2A),
    },
    "bold": {
        "bg": RGBColor(0xF5, 0xF3, 0xFF),
        "title": RGBColor(0x31, 0x0B, 0x8A),
        "body": RGBColor(0x4C, 0x1D, 0x95),
        "muted": RGBColor(0x6D, 0x28, 0xD9),
        "accent": RGBColor(0x7C, 0x3A, 0xED),
        "accent_soft": RGBColor(0xED, 0xE9, 0xFE),
        "card": RGBColor(0xFF, 0xFF, 0xFF),
        "card_border": RGBColor(0xDD, 0xD6, 0xFE),
        "on_accent": RGBColor(0xFF, 0xFF, 0xFF),
    },
    "minimal": {
        "bg": RGBColor(0xFA, 0xFA, 0xFA),
        "title": RGBColor(0x17, 0x17, 0x17),
        "body": RGBColor(0x40, 0x40, 0x40),
        "muted": RGBColor(0x73, 0x73, 0x73),
        "accent": RGBColor(0x17, 0x17, 0x17),
        "accent_soft": RGBColor(0xF0, 0xF0, 0xF0),
        "card": RGBColor(0xFF, 0xFF, 0xFF),
        "card_border": RGBColor(0xE5, 0xE5, 0xE5),
        "on_accent": RGBColor(0xFF, 0xFF, 0xFF),
    },
}


def build_pptx(
    deck: dict[str, Any],
    output_path: str,
    enable_transitions: bool = True,
):
    global CV
    use_template = False

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    # Luôn chuẩn hóa 16:9 để bố cục giống bản không template (tránh chữ tràn khi file mẫu 4:3)
    prs.slide_width = Inches(REF_W)
    prs.slide_height = Inches(REF_H)

    CV = SlideCanvas(w=REF_W, h=REF_H, margin=REF_MARGIN, template_mode=use_template)

    theme = THEMES.get(deck.get("theme", "professional"), THEMES["professional"])
    slides_data = deck.get("slides", [])
    total = len(slides_data)
    deck_title = deck.get("title", "Presentation")

    for slide_index, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        if not CV.template_mode:
            _set_bg(slide, theme["bg"])

        layout = slide_data.get("layout", "bullet")
        content = slide_data.get("content", {})
        image_path = _resolve_slide_image_path(deck, slide_data)

        if layout == "title":
            _render_title(slide, content, theme, image_path)
        elif layout == "bullet":
            _render_bullet(slide, content, theme, image_path)
        elif layout == "two_column":
            _render_two_column(slide, content, theme)
        elif layout == "image_text":
            _render_image_text(slide, content, theme, image_path, slide_data)
        elif layout == "quote":
            _render_quote(slide, content, theme, image_path or _first_deck_image_path(deck))
        elif layout == "closing":
            _render_closing(slide, content, theme)
        else:
            _draw_header_band(slide, theme, content.get("heading", ""))
            _add_text(slide, "Unsupported layout", CV.margin, 1.8, CV.content_w, 0.6, Pt(18), False, theme["body"])

        _draw_footer(slide, theme, slide_index + 1, total, deck_title)

        if slide_data.get("speaker_notes"):
            slide.notes_slide.notes_text_frame.text = slide_data["speaker_notes"]

        if enable_transitions:
            effect = slide_data.get("transition") or ("fade" if slide_index % 2 == 0 else "push")
            _add_slide_transition(slide, effect=effect, speed="med")

    prs.save(output_path)


def _render_title(slide, content, theme, image_path):
    if image_path:
        _add_image_cover(slide, image_path, 0, 0, CV.w, CV.h)
        _add_overlay(slide, 0, 0, CV.w, CV.h, RGBColor(0x00, 0x00, 0x00), transparency=52)
        _add_text(
            slide,
            content.get("heading", ""),
            CV.margin,
            2.15,
            CV.content_w,
            1.35,
            Pt(40),
            True,
            RGBColor(0xFF, 0xFF, 0xFF),
            PP_ALIGN.CENTER,
            max_chars=90,
        )
        _add_text(
            slide,
            content.get("subheading", ""),
            CV.margin,
            3.65,
            CV.content_w,
            0.85,
            Pt(20),
            False,
            RGBColor(0xE2, 0xE8, 0xF0),
            PP_ALIGN.CENTER,
            max_chars=140,
        )
        return

    if not CV.template_mode:
        _add_shape(slide, 0, 0, CV.w, 2.15, theme["accent"], line=False)
        _add_shape(slide, 0, CV.h - 0.32, CV.w, 0.32, theme["accent_soft"], line=False)
        heading_color = RGBColor(0xFF, 0xFF, 0xFF)
    else:
        _add_shape(slide, 0, 0.35, 0.14, 1.05, theme["accent"], line=False)
        heading_color = theme["title"]

    _add_text(
        slide,
        content.get("heading", ""),
        CV.margin,
        0.72 if not CV.template_mode else 0.55,
        CV.content_w,
        1.15,
        Pt(40),
        True,
        heading_color,
        PP_ALIGN.CENTER,
        max_chars=90,
    )
    _add_text(
        slide,
        content.get("subheading", ""),
        CV.margin,
        2.95 if not CV.template_mode else 2.75,
        CV.content_w,
        0.95,
        Pt(20),
        False,
        theme["title"],
        PP_ALIGN.CENTER,
        max_chars=140,
    )


def _render_bullet(slide, content, theme, image_path):
    top = _draw_header_band(slide, theme, content.get("heading", ""))
    bullets = content.get("bullets", [])[:6]
    count = max(len(bullets), 1)
    row_h = min(0.78, 4.55 / count)
    text_left = CV.margin + 0.35
    img_w = 3.55 if image_path else 0
    text_w = CV.content_w - 0.35 - img_w - (0.2 if image_path else 0)

    for i, bullet in enumerate(bullets):
        y = top + 0.12 + i * row_h
        _add_bullet_row(slide, bullet, text_left, y, text_w, row_h - 0.08, theme, Pt(max(14, 19 - count)))

    if image_path:
        _add_image_frame(
            slide,
            image_path,
            CV.w - CV.margin - img_w,
            top + 0.15,
            img_w,
            min(2.35, CV.h - top - 1.05),
            theme,
        )


def _render_two_column(slide, content, theme):
    top = _draw_header_band(slide, theme, content.get("heading", ""))
    col_w = (CV.content_w - 0.35) / 2
    card_h = CV.h - top - 1.05
    if not CV.template_mode:
        _add_card(slide, CV.margin, top + 0.15, col_w, card_h, theme)
        _add_card(slide, CV.margin + col_w + 0.35, top + 0.15, col_w, card_h, theme)
    _add_text(
        slide,
        content.get("left_column", ""),
        CV.margin + 0.2,
        top + 0.35,
        col_w - 0.35,
        card_h - 0.4,
        Pt(15),
        False,
        theme["body"],
        max_chars=420,
    )
    _add_text(
        slide,
        content.get("right_column", ""),
        CV.margin + col_w + 0.55,
        top + 0.35,
        col_w - 0.35,
        card_h - 0.4,
        Pt(15),
        False,
        theme["body"],
        max_chars=420,
    )


def _render_image_text(slide, content, theme, image_path, slide_data):
    top = _draw_header_band(slide, theme, content.get("heading", ""))
    body_h = CV.h - top - 1.05
    text_w = CV.content_w * 0.48
    img_left = CV.margin + text_w + 0.35
    img_w = CV.content_w - text_w - 0.35

    body = content.get("subheading") or content.get("left_column") or _bullets_to_text(content.get("bullets", []))
    _add_text(slide, body, CV.margin, top + 0.2, text_w, body_h, Pt(16), False, theme["body"], max_chars=500)

    if image_path:
        _add_image_frame(slide, image_path, img_left, top + 0.15, img_w, body_h, theme)
    else:
        _add_placeholder(slide, slide_data.get("visual_hint", "Image"), img_left, top + 0.15, img_w, body_h, theme)


def _render_quote(slide, content, theme, image_path):
    if image_path:
        _add_image_cover(slide, image_path, 0, 0, CV.w, CV.h)
        _add_overlay(slide, 0, 0, CV.w, CV.h, RGBColor(0x00, 0x00, 0x00), transparency=55)
        quote_color = RGBColor(0xFF, 0xFF, 0xFF)
        author_color = RGBColor(0xE2, 0xE8, 0xF0)
    else:
        _add_shape(slide, CV.margin, 1.4, 0.12, 3.2, theme["accent"], line=False)
        quote_color = theme["accent"]
        author_color = theme["muted"]

    quote = content.get("quote", "")
    author = content.get("author", "")
    _add_text(
        slide,
        f"\u201c{quote}\u201d",
        CV.margin + 0.5,
        2.0,
        CV.content_w - 0.5,
        2.3,
        Pt(26),
        True,
        quote_color,
        PP_ALIGN.CENTER,
        max_chars=280,
    )
    if author:
        _add_text(
            slide,
            f"\u2014 {author}",
            CV.margin,
            4.55,
            CV.content_w,
            0.55,
            Pt(16),
            False,
            author_color,
            PP_ALIGN.CENTER,
            max_chars=80,
        )


def _render_closing(slide, content, theme):
    if not CV.template_mode:
        _add_shape(slide, 0, 0, CV.w, 0.2, theme["accent"], line=False)
    _add_text(slide, content.get("heading", ""), CV.margin, 2.05, CV.content_w, 1.05, Pt(36), True, theme["title"], PP_ALIGN.CENTER, max_chars=80)
    _add_text(slide, content.get("subheading", ""), CV.margin, 3.15, CV.content_w, 0.75, Pt(19), False, theme["body"], PP_ALIGN.CENTER, max_chars=120)
    cta = content.get("cta", "")
    if cta:
        btn_w = min(5.2, max(2.6, len(cta) * 0.1))
        btn_left = (CV.w - btn_w) / 2
        _add_shape(slide, btn_left, 4.15, btn_w, 0.65, theme["accent"], line=False, radius=True)
        _add_text(slide, cta, btn_left + 0.12, 4.28, btn_w - 0.24, 0.45, Pt(17), True, theme["on_accent"], PP_ALIGN.CENTER, max_chars=60)


def _draw_header_band(slide, theme, heading: str) -> float:
    band_top = 0.35
    band_h = 1.05
    if not CV.template_mode:
        _add_shape(slide, 0, band_top, CV.w, band_h, theme["accent_soft"], line=False)
    _add_shape(slide, 0, band_top, 0.14, band_h, theme["accent"], line=False)
    if heading:
        _add_text(
            slide,
            heading,
            CV.margin + 0.1,
            band_top + 0.16,
            CV.content_w - 0.1,
            0.78,
            Pt(26),
            True,
            theme["title"],
            max_chars=100,
        )
    return band_top + band_h + 0.2


def _draw_footer(slide, theme, page: int, total: int, deck_title: str):
    y = CV.h - 0.42
    _add_shape(slide, CV.margin, y, CV.content_w, 0.02, theme["card_border"], line=False)
    _add_text(slide, deck_title[:48], CV.margin, y + 0.08, CV.content_w * 0.7, 0.28, Pt(9), False, theme["muted"], max_chars=48)
    _add_text(slide, f"{page} / {total}", CV.w - CV.margin - 1.2, y + 0.06, 1.2, 0.3, Pt(9), False, theme["muted"], PP_ALIGN.RIGHT, max_chars=12)


def _add_bullet_row(slide, bullet, left, top, width, height, theme, font_size):
    text = (bullet.get("text") or "").strip()
    detail = (bullet.get("detail") or "").strip()
    _add_text(slide, "\u2022", left, top, 0.22, height, font_size, True, theme["accent"])
    line = f"{text}\n{detail}" if detail else text
    _add_text(slide, line, left + 0.26, top, width - 0.26, height, font_size, False, theme["body"], max_chars=200)


def _add_card(slide, left, top, width, height, theme):
    return _add_shape(slide, left, top, width, height, theme["card"], line=True, border_color=theme["card_border"], radius=True)


def _add_shape(slide, left, top, width, height, fill_color, line=True, border_color=None, radius=False):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line and border_color:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()
    return shape


def _resolve_slide_image_path(deck: dict[str, Any], slide_data: dict[str, Any]) -> str | None:
    content = slide_data.get("content", {})
    direct = content.get("image_path") or slide_data.get("image_path")
    if direct and Path(str(direct)).exists():
        return str(direct)
    image_id = content.get("image_id") or slide_data.get("image_id")
    assets = deck.get("assets", {}) or {}
    if image_id and image_id in assets:
        path = assets[image_id].get("path")
        if path and Path(path).exists():
            return path
    return None


def _first_deck_image_path(deck: dict[str, Any]) -> str | None:
    for meta in (deck.get("assets", {}) or {}).values():
        path = meta.get("path")
        if path and Path(path).exists():
            return path
    return None


def _truncate_text(text: str, max_chars: int | None) -> str:
    if not max_chars or len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip() + "\u2026"


def _add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    size,
    bold,
    color,
    align=PP_ALIGN.LEFT,
    max_chars: int | None = None,
):
    text = _truncate_text(str(text or ""), max_chars)
    tx_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    tf.clear()

    lines = text.splitlines() or [""]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(3)
        p.line_spacing = 1.1
        run = p.add_run()
        run.text = line
        run.font.size = size
        run.font.bold = bold
        run.font.color.rgb = color


def _add_placeholder(slide, text, left, top, width, height, theme):
    shape = _add_shape(slide, left, top, width, height, theme["accent_soft"], line=True, border_color=theme["accent"], radius=True)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.text = _truncate_text(str(text), 80)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].runs[0].font.size = Pt(15)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.color.rgb = theme["accent"]


def _add_image_frame(slide, image_path: str, left, top, width, height, theme):
    _add_shape(slide, left, top, width, height, theme["card"], line=True, border_color=theme["card_border"], radius=True)
    pad = 0.12
    slide.shapes.add_picture(
        image_path,
        Inches(left + pad),
        Inches(top + pad),
        width=Inches(max(0.5, width - 2 * pad)),
        height=Inches(max(0.5, height - 2 * pad)),
    )


def _add_image_cover(slide, image_path: str, left, top, width, height):
    slide.shapes.add_picture(image_path, Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def _add_overlay(slide, left, top, width, height, color, transparency=40):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency
    shape.line.fill.background()


def _bullets_to_text(bullets):
    if not bullets:
        return ""
    return "\n".join(f"• {b.get('text', '')}" for b in bullets[:5])


def _set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_slide_transition(slide, effect: str = "fade", speed: str = "med"):
    sld = slide._element
    for child in list(sld):
        if child.tag.endswith("}transition"):
            sld.remove(child)

    transition = OxmlElement("p:transition")
    transition.set("spd", speed)

    if effect == "push":
        child = OxmlElement("p:push")
        child.set("dir", "l")
    elif effect == "wipe":
        child = OxmlElement("p:wipe")
        child.set("dir", "l")
    else:
        child = OxmlElement("p:fade")

    transition.append(child)

    insert_at = 1
    for idx, child in enumerate(list(sld)):
        if child.tag.endswith("}clrMapOvr"):
            insert_at = idx + 1
    sld.insert(insert_at, transition)
